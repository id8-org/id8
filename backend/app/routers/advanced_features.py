from fastapi import APIRouter, Depends, HTTPException, Body
from sqlalchemy.orm import Session
import logging
from app.services.personalized_idea_service import run_llm_with_user_context
from fastapi.responses import StreamingResponse, JSONResponse
import io
from app.utils.business_utils import BUSINESS_MODEL_GROUPS, BUSINESS_HORIZONTAL_GROUPS, BUSINESS_VERTICAL_GROUPS
from typing import List

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
except ImportError:
    Presentation = None

from ..db import get_db
from ..auth import get_current_user
from ..types import (
    CaseStudy, CaseStudyCreate, CaseStudyRequest,
    MarketSnapshot, MarketSnapshotCreate, MarketSnapshotRequest,
    LensInsight, LensInsightCreate, LensInsightRequest,
    VCThesisComparison, VCThesisComparisonCreate, VCThesisComparisonRequest,
    InvestorDeck, InvestorDeckCreate, InvestorDeckRequest,
    ExportRecordCreate, ExportRecordOut
)
from app.models import (
    CaseStudy as CaseStudyModel,
    MarketSnapshot as MarketSnapshotModel,
    LensInsight as LensInsightModel,
    VCThesisComparison as VCThesisComparisonModel,
    InvestorDeck as InvestorDeckModel,
    Idea, ExportRecord, DeepDiveVersion
)
from app.llm_center.legacy_wrappers import (
    generate_case_study,
    generate_market_snapshot,
    generate_lens_insight,
    generate_vc_thesis_comparison,
    generate_investor_deck
)
import uuid
from app.types import DeepDiveCreate, DeepDiveOut
from app.services.suggested_service import (
    save_suggested, get_suggested_by_id, get_suggested_by_idea_id, update_suggested
)
from app.services.iterating_service import (
    create_iterating, get_iterating_by_id, get_latest_iterating_by_idea, update_iterating, delete_iterating
)
from app.types import SuggestedCreate, SuggestedOut, IteratingCreate, IteratingOut

logger = logging.getLogger(__name__)
router = APIRouter(tags=["advanced"])

def get_all_idea_context(db: Session, idea_id: str) -> dict[str, object]:
    """Fetch all advanced feature data for an idea and return as a dict."""
    idea = db.query(Idea).filter(Idea.id == idea_id).first()
    deep_dive = idea.deep_dive if idea and idea.deep_dive else {}
    case_study = db.query(CaseStudyModel).filter(CaseStudyModel.idea_id == idea_id).first()
    market_snapshot = db.query(MarketSnapshotModel).filter(MarketSnapshotModel.idea_id == idea_id).first()
    lens_insights = db.query(LensInsightModel).filter(LensInsightModel.idea_id == idea_id).all()
    vc_thesis = db.query(VCThesisComparisonModel).filter(VCThesisComparisonModel.idea_id == idea_id).all()
    investor_deck = db.query(InvestorDeckModel).filter(InvestorDeckModel.idea_id == idea_id).first()
    return {
        'deep_dive': deep_dive,
        'case_study': case_study.llm_raw_response if case_study else None,
        'market_snapshot': market_snapshot.llm_raw_response if market_snapshot else None,
        'lens_insights': [li.llm_raw_response for li in lens_insights] if lens_insights else [],
        'vc_thesis': [vc.llm_raw_response for vc in vc_thesis] if vc_thesis else [],
        'investor_deck': investor_deck.llm_raw_response if investor_deck else None
    }

@router.post("/case-study")
async def create_case_study(
    request: CaseStudyRequest,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Generate a case study for an idea."""
    try:
        idea = db.query(Idea).filter(Idea.id == request.idea_id).first()
        if idea is None:
            raise HTTPException(status_code=404, detail="Idea not found")
        existing_case_study = db.query(CaseStudyModel).filter(
            CaseStudyModel.idea_id == request.idea_id
        ).first()
        if existing_case_study is not None:
            return {
                "case_study": CaseStudy.model_validate(existing_case_study),
                "llm_raw_response": existing_case_study.llm_raw_response
            }
        all_context = get_all_idea_context(db, request.idea_id)
        idea_data = {
            'title': idea.title,
            'hook': idea.hook,
            'value': idea.value,
            'evidence': idea.evidence,
            'differentiator': idea.differentiator,
            'all_context': all_context
        }
        llm_response = await run_llm_with_user_context(
            user=current_user,
            db=db,
            llm_func=generate_case_study,
            idea_data=idea_data,
            extra_args={"company_name": request.company_name}
        )
        case_study_data = CaseStudyCreate(
            company_name=llm_response.get('company_name', 'Unknown Company'),
            industry=llm_response.get('industry'),
            business_model=llm_response.get('business_model'),
            success_factors=llm_response.get('success_factors'),
            challenges=llm_response.get('challenges'),
            lessons_learned=llm_response.get('lessons_learned'),
            market_size=llm_response.get('market_size'),
            funding_raised=llm_response.get('funding_raised'),
            exit_value=llm_response.get('exit_value')
        )
        case_study = CaseStudyModel(
            idea_id=request.idea_id,
            llm_raw_response=str(llm_response),
            **case_study_data.dict()
        )
        db.add(case_study)
        db.commit()
        db.refresh(case_study)
        return {
            "case_study": CaseStudy.model_validate(case_study),
            "llm_raw_response": str(llm_response)
        }
    except Exception as e:
        logger.error(f"Error creating case study: {e}")
        raise HTTPException(status_code=500, detail="Failed to generate case study")

@router.get("/case-study/{idea_id}")
async def get_case_study(
    idea_id: str,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Get case study for an idea."""
    case_study = db.query(CaseStudyModel).filter(
        CaseStudyModel.idea_id == idea_id
    ).first()
    
    if case_study is None:
        raise HTTPException(status_code=404, detail="Case study not found")
    
    return {
        "case_study": CaseStudy.model_validate(case_study),
        "llm_raw_response": case_study.llm_raw_response
    }

@router.post("/market-snapshot")
async def create_market_snapshot(
    request: MarketSnapshotRequest,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Generate a market snapshot for an idea."""
    try:
        idea = db.query(Idea).filter(Idea.id == request.idea_id).first()
        if idea is None:
            raise HTTPException(status_code=404, detail="Idea not found")
        existing_snapshot = db.query(MarketSnapshotModel).filter(
            MarketSnapshotModel.idea_id == request.idea_id
        ).first()
        if existing_snapshot is not None:
            return {
                "market_snapshot": MarketSnapshot.model_validate(existing_snapshot),
                "llm_raw_response": existing_snapshot.llm_raw_response
            }
        all_context = get_all_idea_context(db, request.idea_id)
        idea_data = {
            'title': idea.title,
            'hook': idea.hook,
            'value': idea.value,
            'evidence': idea.evidence,
            'differentiator': idea.differentiator,
            'all_context': all_context
        }
        llm_response = await run_llm_with_user_context(
            user=current_user,
            db=db,
            llm_func=generate_market_snapshot,
            idea_data=idea_data
        )
        snapshot_data = MarketSnapshotCreate(
            market_size=llm_response.get('market_size'),
            growth_rate=llm_response.get('growth_rate'),
            key_players=llm_response.get('key_players', []),
            market_trends=llm_response.get('market_trends'),
            regulatory_environment=llm_response.get('regulatory_environment'),
            competitive_landscape=llm_response.get('competitive_landscape'),
            entry_barriers=llm_response.get('entry_barriers')
        )
        snapshot = MarketSnapshotModel(
            idea_id=request.idea_id,
            llm_raw_response=str(llm_response),
            **snapshot_data.dict()
        )
        db.add(snapshot)
        db.commit()
        db.refresh(snapshot)
        return {
            "market_snapshot": MarketSnapshot.model_validate(snapshot),
            "llm_raw_response": str(llm_response)
        }
    except Exception as e:
        logger.error(f"Error creating market snapshot: {e}")
        raise HTTPException(status_code=500, detail="Failed to generate market snapshot")

@router.get("/market-snapshot/{idea_id}")
async def get_market_snapshot(
    idea_id: str,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Get market snapshot for an idea."""
    snapshot = db.query(MarketSnapshotModel).filter(
        MarketSnapshotModel.idea_id == idea_id
    ).first()
    
    if snapshot is None:
        raise HTTPException(status_code=404, detail="Market snapshot not found")
    
    return {
        "market_snapshot": MarketSnapshot.model_validate(snapshot),
        "llm_raw_response": snapshot.llm_raw_response
    }

@router.post("/lens-insight")
async def create_lens_insight(
    request: LensInsightRequest,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Generate lens insights for an idea."""
    try:
        idea = db.query(Idea).filter(Idea.id == request.idea_id).first()
        if idea is None:
            raise HTTPException(status_code=404, detail="Idea not found")
        existing_insight = db.query(LensInsightModel).filter(
            LensInsightModel.idea_id == request.idea_id,
            LensInsightModel.lens_type == request.lens_type
        ).first()
        if existing_insight is not None:
            return {
                "lens_insight": LensInsight.model_validate(existing_insight),
                "llm_raw_response": existing_insight.llm_raw_response
            }
        all_context = get_all_idea_context(db, request.idea_id)
        idea_data = {
            'title': idea.title,
            'hook': idea.hook,
            'value': idea.value,
            'evidence': idea.evidence,
            'differentiator': idea.differentiator,
            'all_context': all_context
        }
        llm_response = await run_llm_with_user_context(
            user=current_user,
            db=db,
            llm_func=generate_lens_insight,
            idea_data=idea_data,
            extra_args={"lens_type": request.lens_type}
        )
        insight_data = LensInsightCreate(
            lens_type=request.lens_type,
            insights=llm_response.get('insights'),
            opportunities=llm_response.get('opportunities'),
            risks=llm_response.get('risks'),
            recommendations=llm_response.get('recommendations')
        )
        insight = LensInsightModel(
            idea_id=request.idea_id,
            llm_raw_response=str(llm_response),
            **insight_data.dict()
        )
        db.add(insight)
        db.commit()
        db.refresh(insight)
        return {
            "lens_insight": LensInsight.model_validate(insight),
            "llm_raw_response": str(llm_response)
        }
    except Exception as e:
        logger.error(f"Error creating lens insight: {e}")
        raise HTTPException(status_code=500, detail="Failed to generate lens insight")

@router.get("/lens-insights/{idea_id}")
async def get_lens_insights(
    idea_id: str,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Get all lens insights for an idea."""
    insights = db.query(LensInsightModel).filter(
        LensInsightModel.idea_id == idea_id
    ).all()
    
    if insights is None or len(insights) == 0:
        return {"lens_insights": [], "llm_raw_responses": {}}
    
    return {
        "lens_insights": [LensInsight.model_validate(insight) for insight in insights],
        "llm_raw_responses": {insight.lens_type: insight.llm_raw_response for insight in insights}
    }

@router.post("/vc-thesis-comparison")
async def create_vc_thesis_comparison(
    request: VCThesisComparisonRequest,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Generate VC thesis comparison for an idea."""
    try:
        idea = db.query(Idea).filter(Idea.id == request.idea_id).first()
        if idea is None:
            raise HTTPException(status_code=404, detail="Idea not found")
        if request.vc_firm:
            existing_comparison = db.query(VCThesisComparisonModel).filter(
                VCThesisComparisonModel.idea_id == request.idea_id,
                VCThesisComparisonModel.vc_firm == request.vc_firm
            ).first()
            if existing_comparison is not None:
                return {
                    "vc_thesis_comparison": VCThesisComparison.model_validate(existing_comparison),
                    "llm_raw_response": existing_comparison.llm_raw_response
                }
        all_context = get_all_idea_context(db, request.idea_id)
        idea_data = {
            'title': idea.title,
            'hook': idea.hook,
            'value': idea.value,
            'evidence': idea.evidence,
            'differentiator': idea.differentiator,
            'all_context': all_context
        }
        llm_response = await run_llm_with_user_context(
            user=current_user,
            db=db,
            llm_func=generate_vc_thesis_comparison,
            idea_data=idea_data,
            extra_args={"vc_firm": request.vc_firm}
        )
        comparison_data = VCThesisComparisonCreate(
            vc_firm=llm_response.get('vc_firm', 'Unknown VC'),
            thesis_focus=llm_response.get('thesis_focus'),
            alignment_score=llm_response.get('alignment_score'),
            key_alignment_points=llm_response.get('key_alignment_points'),
            potential_concerns=llm_response.get('potential_concerns'),
            investment_likelihood=llm_response.get('investment_likelihood')
        )
        comparison = VCThesisComparisonModel(
            idea_id=request.idea_id,
            llm_raw_response=str(llm_response),
            **comparison_data.dict()
        )
        db.add(comparison)
        db.commit()
        db.refresh(comparison)
        return {
            "vc_thesis_comparison": VCThesisComparison.model_validate(comparison),
            "llm_raw_response": str(llm_response)
        }
    except Exception as e:
        logger.error(f"Error creating VC thesis comparison: {e}")
        raise HTTPException(status_code=500, detail="Failed to generate VC thesis comparison")

@router.get("/vc-thesis-comparisons/{idea_id}")
async def get_vc_thesis_comparisons(
    idea_id: str,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Get all VC thesis comparisons for an idea."""
    comparisons = db.query(VCThesisComparisonModel).filter(
        VCThesisComparisonModel.idea_id == idea_id
    ).all()
    
    if comparisons is None or len(comparisons) == 0:
        return {"vc_thesis_comparisons": [], "llm_raw_responses": {}}
    
    return {
        "vc_thesis_comparisons": [VCThesisComparison.model_validate(comparison) for comparison in comparisons],
        "llm_raw_responses": {comparison.vc_firm: comparison.llm_raw_response for comparison in comparisons}
    }

@router.post("/investor-deck")
async def create_investor_deck(
    request: InvestorDeckRequest,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Generate an investor deck for an idea."""
    try:
        idea = db.query(Idea).filter(Idea.id == request.idea_id).first()
        if idea is None:
            raise HTTPException(status_code=404, detail="Idea not found")
        # Query for existing deck with same options
        existing_deck = db.query(InvestorDeckModel).filter(
            InvestorDeckModel.idea_id == request.idea_id,
            InvestorDeckModel.focus_area == request.focus_area,
            InvestorDeckModel.style == request.style,
            InvestorDeckModel.slides == request.slides
        ).first()
        if existing_deck is not None:
            return {
                "investor_deck": InvestorDeck.model_validate(existing_deck),
                "llm_raw_response": existing_deck.llm_raw_response
            }
        all_context = get_all_idea_context(db, request.idea_id)
        idea_data = {
            'title': idea.title,
            'hook': idea.hook,
            'value': idea.value,
            'evidence': idea.evidence,
            'differentiator': idea.differentiator,
            'all_context': all_context
        }
        llm_response = await run_llm_with_user_context(
            user=current_user,
            db=db,
            llm_func=generate_investor_deck,
            idea_data=idea_data,
            extra_args={
                "include_case_studies": request.include_case_studies,
                "include_market_analysis": request.include_market_analysis,
                "include_financial_projections": request.include_financial_projections,
                "focus_area": request.focus_area,
                "slides": request.slides,
                "style": request.style
            }
        )
        deck_data = InvestorDeckCreate(
            deck_content=llm_response,
            slides=request.slides,
            focus_area=request.focus_area,
            style=request.style
        )
        deck = InvestorDeckModel(
            idea_id=request.idea_id,
            llm_raw_response=str(llm_response),
            slides=request.slides,
            focus_area=request.focus_area,
            style=request.style,
            **deck_data.dict()
        )
        db.add(deck)
        db.commit()
        db.refresh(deck)
        return {
            "investor_deck": InvestorDeck.model_validate(deck),
            "llm_raw_response": str(llm_response)
        }
    except Exception as e:
        logger.error(f"Error creating investor deck: {e}")
        raise HTTPException(status_code=500, detail="Failed to generate investor deck")

@router.get("/investor-deck/{idea_id}")
async def get_investor_deck(
    idea_id: str,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Get investor deck for an idea."""
    deck = db.query(InvestorDeckModel).filter(
        InvestorDeckModel.idea_id == idea_id
    ).first()
    
    if deck is None:
        raise HTTPException(status_code=404, detail="Investor deck not found")
    
    return {
        "investor_deck": InvestorDeck.model_validate(deck),
        "llm_raw_response": deck.llm_raw_response
    }

@router.get("/investor-deck/{idea_id}/pptx")
def export_investor_deck_pptx(
    idea_id: str,
    slides: str = "",
    focus_area: str = "",
    style: str = "modern",
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Export the investor deck as a PowerPoint (PPTX) file."""
    if Presentation is None:
        return HTTPException(status_code=501, detail="python-pptx is not installed on the server.")
    # Parse slides param if provided
    slides_list = None
    if slides:
        import json
        try:
            slides_list = json.loads(slides)
        except Exception:
            slides_list = None
    # Query for deck with matching options if provided
    query = db.query(InvestorDeckModel).filter(InvestorDeckModel.idea_id == idea_id)
    if focus_area and focus_area != "":
        query = query.filter(InvestorDeckModel.focus_area == focus_area)
    if style and style != "":
        query = query.filter(InvestorDeckModel.style == style)
    if slides_list is not None:
        query = query.filter(InvestorDeckModel.slides == slides_list)
    deck = query.first()
    if deck is None or not deck.deck_content:
        raise HTTPException(status_code=404, detail="Investor deck not found or not generated yet.")
    deck_data = deck.deck_content
    prs = Presentation()
    # Title slide
    title = deck_data.get('title', 'Investor Deck')
    slides_content = deck_data.get('slides', [])
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated by ID8"
    # Content slides
    for s in slides_content:
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s.get('title', 'Slide')
        content = s.get('content') or s.get('key_points') or []
        if isinstance(content, str):
            content = [content]
        body = slide.placeholders[1].text_frame
        for para in content:
            p = body.add_paragraph()
            p.text = para
            p.font.size = Pt(18)
    # Output to bytes
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    filename = f"investor_deck_{idea_id}.pptx"
    return StreamingResponse(
        pptx_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

@router.post("/exports", response_model=ExportRecordOut)
def create_export_record(
    record: ExportRecordCreate,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Log a new export record (called after deck generation/download)."""
    export = ExportRecord(
        user_id=current_user.id,
        idea_id=record.idea_id,
        deck_id=record.deck_id,
        slides=record.slides,
        focus_area=record.focus_area,
        style=record.style,
        file_type=record.file_type,
        recipient=record.recipient
    )
    db.add(export)
    db.commit()
    db.refresh(export)
    return export

@router.get("/exports", response_model=List[ExportRecordOut])
def list_export_records(
    idea_id: str = "",
    recipient: str = "",
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """List all export records for the current user (optionally filter by idea_id, recipient, etc.)."""
    query = db.query(ExportRecord).filter(ExportRecord.user_id == current_user.id)
    if idea_id and idea_id != "":
        query = query.filter(ExportRecord.idea_id == idea_id)
    if recipient and recipient != "":
        query = query.filter(ExportRecord.recipient == recipient)
    return query.order_by(ExportRecord.created_at.desc()).all()

@router.get("/exports/{export_id}", response_model=ExportRecordOut)
def get_export_record(
    export_id: str,
    db: Session = Depends(get_db),
    current_user = Depends(get_current_user)
):
    """Get details for a specific export record."""
    export = db.query(ExportRecord).filter(ExportRecord.id == export_id, ExportRecord.user_id == current_user.id).first()
    if export is None:
        raise HTTPException(status_code=404, detail="Export record not found")
    return export

@router.get("/meta/business-options", tags=["meta"])
def get_business_options():
    """Return all business model, horizontal, and vertical groups for onboarding/profile/ideas."""
    return JSONResponse({
        "business_model_groups": BUSINESS_MODEL_GROUPS,
        "business_horizontal_groups": BUSINESS_HORIZONTAL_GROUPS,
        "business_vertical_groups": BUSINESS_VERTICAL_GROUPS,
    }) 

# Deep dive and iterating endpoints are now part of the main advanced features router

@router.post("/", response_model=DeepDiveOut)
def create_deep_dive(deep_dive: DeepDiveCreate, db: Session = Depends(get_db)):
    idea = db.query(Idea).filter(Idea.id == deep_dive.idea_id).first()
    if idea is None:
        raise HTTPException(status_code=404, detail="Idea not found")
    new_deep_dive = DeepDiveVersion(
        id=str(uuid.uuid4()),
        idea_id=deep_dive.idea_id,
        version_number=deep_dive.version,  # type: ignore[assignment]
        fields=deep_dive.data,  # type: ignore[assignment]
        llm_raw_response=deep_dive.llm_raw_response or ""  # type: ignore[assignment]
    )
    db.add(new_deep_dive)
    db.commit()
    db.refresh(new_deep_dive)
    return new_deep_dive

@router.get("/{deep_dive_id}", response_model=DeepDiveOut)
def get_deep_dive(deep_dive_id: str, db: Session = Depends(get_db)):
    deep_dive = db.query(DeepDiveVersion).filter(DeepDiveVersion.id == deep_dive_id).first()
    if deep_dive is None:
        raise HTTPException(status_code=404, detail="Deep Dive not found")
    return deep_dive

@router.get("/idea/{idea_id}", response_model=DeepDiveOut)
def get_deep_dive_by_idea(idea_id: str, db: Session = Depends(get_db)):
    deep_dive = db.query(DeepDiveVersion).filter(DeepDiveVersion.idea_id == idea_id).order_by(DeepDiveVersion.version_number.desc()).first()
    if deep_dive is None:
        raise HTTPException(status_code=404, detail="Deep Dive not found for idea")
    return deep_dive

@router.put("/{deep_dive_id}", response_model=DeepDiveOut)
def update_deep_dive(deep_dive_id: str, deep_dive: DeepDiveCreate, db: Session = Depends(get_db)):
    db_deep_dive = db.query(DeepDiveVersion).filter(DeepDiveVersion.id == deep_dive_id).first()
    if db_deep_dive is None:
        raise HTTPException(status_code=404, detail="Deep Dive not found")
    db_deep_dive.fields = deep_dive.data  # type: ignore[assignment]
    db_deep_dive.version_number = deep_dive.version  # type: ignore[assignment]
    db_deep_dive.llm_raw_response = deep_dive.llm_raw_response or ""  # type: ignore[assignment]
    db.commit()
    db.refresh(db_deep_dive)
    return db_deep_dive

@router.delete("/{deep_dive_id}")
def delete_deep_dive(deep_dive_id: str, db: Session = Depends(get_db)):
    db_deep_dive = db.query(DeepDiveVersion).filter(DeepDiveVersion.id == deep_dive_id).first()
    if db_deep_dive is None:
        raise HTTPException(status_code=404, detail="Deep Dive not found")
    db.delete(db_deep_dive)
    db.commit()
    return {"detail": "Deep Dive deleted"} 

# --- Suggested Endpoints ---
@router.post("/suggested/", response_model=SuggestedOut)
def create_suggested_api(suggested: SuggestedCreate, db: Session = Depends(get_db)):
    return save_suggested(db, int(suggested.idea_id), suggested.data, suggested.version, suggested.llm_raw_response)

@router.get("/suggested/{suggested_id}", response_model=SuggestedOut)
def get_suggested_api(suggested_id: str, db: Session = Depends(get_db)):
    result = get_suggested_by_id(db, int(suggested_id))
    if not result:
        raise HTTPException(status_code=404, detail="Suggested not found")
    return result

@router.get("/suggested/idea/{idea_id}", response_model=SuggestedOut)
def get_latest_suggested_by_idea_api(idea_id: str, db: Session = Depends(get_db)):
    result = get_suggested_by_idea_id(db, int(idea_id))
    if not result:
        raise HTTPException(status_code=404, detail="Suggested not found for idea")
    return result

@router.put("/suggested/{suggested_id}", response_model=SuggestedOut)
def update_suggested_api(suggested_id: str, suggested: SuggestedCreate, db: Session = Depends(get_db)):
    result = update_suggested(db, int(suggested_id), suggested.data, suggested.version, suggested.llm_raw_response)
    if not result:
        raise HTTPException(status_code=404, detail="Suggested not found")
    return result

# --- Iterating Endpoints ---
@router.post("/iterating/", response_model=IteratingOut)
def create_iterating_api(iterating: IteratingCreate, db: Session = Depends(get_db)):
    return create_iterating(db, iterating.idea_id, iterating.data, iterating.version, iterating.llm_raw_response)

@router.get("/iterating/{iterating_id}", response_model=IteratingOut)
def get_iterating_api(iterating_id: str, db: Session = Depends(get_db)):
    result = get_iterating_by_id(db, iterating_id)
    if not result:
        raise HTTPException(status_code=404, detail="Iterating not found")
    return result

@router.get("/iterating/idea/{idea_id}", response_model=IteratingOut)
def get_latest_iterating_by_idea_api(idea_id: str, db: Session = Depends(get_db)):
    result = get_latest_iterating_by_idea(db, idea_id)
    if not result:
        raise HTTPException(status_code=404, detail="Iterating not found for idea")
    return result

@router.put("/iterating/{iterating_id}", response_model=IteratingOut)
def update_iterating_api(iterating_id: str, iterating: IteratingCreate, db: Session = Depends(get_db)):
    result = update_iterating(db, iterating_id, iterating.data, iterating.version, iterating.llm_raw_response)
    if not result:
        raise HTTPException(status_code=404, detail="Iterating not found")
    return result

@router.delete("/iterating/{iterating_id}")
def delete_iterating_api(iterating_id: str, db: Session = Depends(get_db)):
    ok = delete_iterating(db, iterating_id)
    if not ok:
        raise HTTPException(status_code=404, detail="Iterating not found")
    return {"detail": "Iterating deleted"} 